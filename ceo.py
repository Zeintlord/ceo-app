import os
import requests
import pandas as pd
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation

# =====================================
# CONFIGURACIÓN INICIAL
# =====================================

# API key de Google Gemini (incluida directamente)
API_KEY = "AIzaSyD7iuusLHEGLOXAGfH42wENZ5ujbekozJI"

# Límite de caracteres para el contenido de cada agente
MAX_CHARS = 10000

# Endpoint de la API de Gemini
API_ENDPOINT = "https://generativelanguage.googleapis.com/v1beta/models/gemini-2.0-flash:generateContent"

# Mensajes de sistema para cada uno de los 10 agentes
MENSAJE_SISTEMA_SALES = (
    "Eres un experto en ventas B2B especializado en prospección outbound y cierre de reuniones comerciales. "
    "Tu rol es optimizar procesos comerciales con foco en conversiones, análisis del ICP y eficiencia en campañas multicanal. "
    "Tienes acceso a secuencias de email, CRM, estrategias de cold outreach y documentos como pitch de ventas, análisis de competencia y flujos de seguimiento. "
    "Contexto de Saint Analytics: Agencia de automatización de reportes y data analytics, con enfoque en empresas que tercerizan inteligencia comercial. "
    "Documentos disponibles:\n[LISTA_DOCUMENTOS]\n\nContenido:\n[contenido de archivos]"
)

MENSAJE_SISTEMA_MARKETING = (
    "Eres un estratega en marketing digital orientado a generación de leads B2B mediante contenido, branding y automatización. "
    "Diseñas campañas orgánicas y pagas en LinkedIn, gestionas lead magnets, pitch deck, e-books y casos de éxito. "
    "Tu enfoque está en el posicionamiento de marca, contenido educativo, SEO y adquisición multicanal. "
    "Contexto: Saint Analytics ofrece soluciones analíticas personalizadas con fuerte presencia en redes y uso intensivo de IA para personalización. "
    "Documentos disponibles:\n[LISTA_DOCUMENTOS]\n\nContenido:\n[contenido de archivos]"
)

MENSAJE_SISTEMA_TOOLS = (
    "Eres un especialista en implementación de herramientas tecnológicas y automatización de flujos de trabajo. "
    "Diseñas e integras software como Instantly, Sales Navigator, ProspEO, Zapier, Google Drive, CRM y sistemas de scraping. "
    "Optimiza procesos comerciales y operativos con bajo costo y alta escalabilidad. "
    "Tu conocimiento es clave para potenciar las capacidades de Saint Analytics en generación de leads y eficiencia interna. "
    "Documentos disponibles:\n[LISTA_DOCUMENTOS]\n\nContenido:\n[contenido de archivos]"
)

MENSAJE_SISTEMA_TASKS = (
    "Eres un gestor de proyectos orientado a la eficiencia, priorización y cumplimiento de entregables. "
    "Supervisas tareas por área (ventas, marketing, branding, legales, etc.), optimizas flujos de trabajo y aseguras la ejecución de iniciativas clave. "
    "Tu enfoque es garantizar que todas las tareas estén alineadas con la estrategia de crecimiento de Saint Analytics. "
    "Documentos disponibles:\n[LISTA_DOCUMENTOS]\n\nContenido:\n[contenido de archivos]"
)

MENSAJE_SISTEMA_BRANDING = (
    "Eres un experto en identidad de marca y comunicación visual. "
    "Te encargas del desarrollo del manual de marca, logotipos, presentaciones y activos gráficos coherentes con el posicionamiento B2B de Saint Analytics. "
    "Tu rol es elevar la percepción de valor de la empresa mediante coherencia visual y emocional. "
    "Documentos disponibles:\n[LISTA_DOCUMENTOS]\n\nContenido:\n[contenido de archivos]"
)

MENSAJE_SISTEMA_FINANCE = (
    "Eres un asesor financiero enfocado en planificación, proyecciones y estructura de precios para servicios B2B. "
    "Manejas presupuestos, flujos de ingreso, estrategias de cobro internacional (Stripe, PayPal) y herramientas de control de gastos. "
    "Tu responsabilidad es asegurar la sostenibilidad financiera y la escalabilidad de Saint Analytics en mercados locales e internacionales. "
    "Documentos disponibles:\n[LISTA_DOCUMENTOS]\n\nContenido:\n[contenido de archivos]"
)

MENSAJE_SISTEMA_LEGALES = (
    "Eres un abogado corporativo especializado en startups y servicios digitales. "
    "Te encargas de contratos, NDAs, términos y condiciones, pactos de socios y la estructura legal-operativa de Saint Analytics (Argentina/EE.UU). "
    "Tu prioridad es garantizar cumplimiento normativo, seguridad jurídica y escalabilidad. "
    "Documentos disponibles:\n[LISTA_DOCUMENTOS]\n\nContenido:\n[contenido de archivos]"
)

MENSAJE_SISTEMA_CLIENTES = (
    "Eres un experto en gestión de relaciones con clientes (Customer Success y CRM). "
    "Supervisas onboarding, SLA, segmentación de clientes (PYMEs, freelance, EE.UU.) y experiencia posventa. "
    "Tu misión es aumentar la retención, generar testimonios y asegurar relaciones de largo plazo para Saint Analytics. "
    "Documentos disponibles:\n[LISTA_DOCUMENTOS]\n\nContenido:\n[contenido de archivos]"
)

MENSAJE_SISTEMA_WEB = (
    "Eres un desarrollador web con enfoque en performance y conversión. "
    "Gestionas la página de Saint Analytics, formularios, validadores, SEO técnico y funcionalidades como calculadora ROI, blogs automáticos y contacto inmediato. "
    "Tu objetivo es maximizar la captación digital y reforzar la confianza de los visitantes. "
    "Documentos disponibles:\n[LISTA_DOCUMENTOS]\n\nContenido:\n[contenido de archivos]"
)

MENSAJE_SISTEMA_HR = (
    "Eres un responsable de Recursos Humanos enfocado en reclutamiento y crecimiento del equipo. "
    "Definís perfiles, filtras candidatos, coordinas entrevistas y manejas documentación laboral. "
    "Aseguras que Saint Analytics cuente con el talento necesario para escalar con agilidad. "
    "Documentos disponibles:\n[LISTA_DOCUMENTOS]\n\nContenido:\n[contenido de archivos]"
)

MENSAJE_SISTEMA_CEO = (
    "Eres el CEO de Saint Analytics. Tomas decisiones estratégicas basadas en las recomendaciones de expertos en ventas, marketing, finanzas, legales y operaciones. "
    "Tu rol es coordinar, priorizar y maximizar el impacto organizacional, alineando cada acción con la visión de escalar soluciones data-driven automatizadas en el mercado global B2B."
)

# Diccionario que mapea cada agente con su carpeta local y su mensaje de sistema
AGENTES = {
    "Sales":              ("contexto Sales",               MENSAJE_SISTEMA_SALES),
    "Marketing":          ("contexto Marketing",           MENSAJE_SISTEMA_MARKETING),
    "Tools":              ("contexto Tools",               MENSAJE_SISTEMA_TOOLS),
    "Tasks Pendientes":   ("contexto TasksPendientes",     MENSAJE_SISTEMA_TASKS),
    "Branding":           ("contexto Branding",            MENSAJE_SISTEMA_BRANDING),
    "Finance":            ("contexto Finance",             MENSAJE_SISTEMA_FINANCE),
    "Legales & Documentos": ("contexto LegalesDocumentos", MENSAJE_SISTEMA_LEGALES),
    "Clientes":           ("contexto Clientes",            MENSAJE_SISTEMA_CLIENTES),
    "Web":                ("contexto Web",                 MENSAJE_SISTEMA_WEB),
    "HR":                 ("contexto HR",                  MENSAJE_SISTEMA_HR)
}

# =====================================
# FUNCIONES PARA LECTURA DE ARCHIVOS
# =====================================

def parse_file(ruta_archivo):
    """
    Devuelve el texto extraído del archivo y el nombre del archivo.
    Soporta PDF, DOC, DOCX, XLS, XLSX, CSV, TXT, PPTX.
    Ignora otros formatos.
    """
    extension = ruta_archivo.lower().split('.')[-1]
    nombre_archivo = os.path.basename(ruta_archivo)
    texto = ""

    try:
        if extension == 'pdf':
            with open(ruta_archivo, 'rb') as f:
                reader = PdfReader(f)
                for page in reader.pages:
                    page_text = page.extract_text() or ""
                    texto += page_text + "\n"
        elif extension in ['doc', 'docx']:
            doc = Document(ruta_archivo)
            for para in doc.paragraphs:
                texto += para.text + "\n"
        elif extension in ['xls', 'xlsx']:
            df = pd.read_excel(ruta_archivo)
            texto += df.to_string() + "\n"
        elif extension == 'csv':
            df = pd.read_csv(ruta_archivo)
            texto += df.to_string() + "\n"
        elif extension == 'txt':
            with open(ruta_archivo, 'r', encoding='utf-8') as f:
                texto += f.read() + "\n"
        elif extension == 'pptx':
            pres = Presentation(ruta_archivo)
            for slide in pres.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        texto += shape.text + "\n"
    except Exception as e:
        print(f"Error leyendo '{ruta_archivo}': {e}")

    return texto, nombre_archivo

def leer_carpeta_recursiva(carpeta):
    """
    Explora recursivamente la carpeta y sus subcarpetas, 
    parsea los archivos soportados y concatena su texto, 
    y guarda la lista de nombres de archivo.
    Retorna (texto_concatenado, lista_nombres).
    """
    contenido_total = ""
    lista_nombres = []
    if not os.path.exists(carpeta) or not os.path.isdir(carpeta):
        return "", []

    for root, dirs, files in os.walk(carpeta):
        for archivo in files:
            ruta_archivo = os.path.join(root, archivo)
            texto, nombre = parse_file(ruta_archivo)
            if texto:
                lista_nombres.append(nombre)
                contenido_total += texto
                if len(contenido_total) >= MAX_CHARS:
                    return contenido_total[:MAX_CHARS], lista_nombres

    return contenido_total[:MAX_CHARS], lista_nombres

# =====================================
# FUNCIONES PARA LA API
# =====================================

def obtener_respuesta_api(prompt):
    """
    Envía el prompt a la API de Google Gemini y devuelve la respuesta generada.
    """
    url = f"{API_ENDPOINT}?key={API_KEY}"
    data = {
        "contents": [{
            "parts": [{"text": prompt}]
        }]
    }
    try:
        response = requests.post(url, json=data)
        response.raise_for_status()
        return response.json()['candidates'][0]['content']['parts'][0]['text']
    except Exception as e:
        print(f"Error en la llamada a la API: {e}")
        return None

def construir_prompt_agente(mensaje_sistema, contenido_contexto, lista_docs, solicitud_usuario):
    """
    Construye el prompt para un agente:
      - Reemplaza [LISTA_DOCUMENTOS] con la lista de archivos.
      - Reemplaza [contenido de archivos] con el texto extraído.
      - Añade la solicitud del usuario.
    """
    doc_string = "\n".join(lista_docs) if lista_docs else "Ningún archivo."
    prompt = mensaje_sistema.replace("[LISTA_DOCUMENTOS]", doc_string)
    prompt = prompt.replace("[contenido de archivos]", contenido_contexto)
    prompt += f"\n\nSolicitud del usuario: {solicitud_usuario}"
    return prompt

def construir_prompt_ceo(respuestas_agentes):
    """
    Construye el prompt para el agente CEO, integrando las respuestas de todos los agentes.
    """
    prompt = MENSAJE_SISTEMA_CEO + "\n\nAquí están las respuestas de los agentes:\n"
    for area, respuesta in respuestas_agentes.items():
        prompt += f"{area}:\n{respuesta}\n\n"
    prompt += "Basado en esta información, por favor proporciona una decisión o conclusión final."
    return prompt

# =====================================
# FUNCIÓN PRINCIPAL CON BUCLE DE PREGUNTAS
# =====================================

def main():
    while True:
        solicitud_usuario = input("\nIngrese su solicitud (o 'salir' para terminar): ").strip()
        if solicitud_usuario.lower() == "salir":
            print("Saliendo del sistema...")
            break

        # Permite opcionalmente procesar un solo agente:
        area_input = input("Ingrese el área a procesar (deje vacío para todas): ").strip()

        if area_input:
            if area_input in AGENTES:
                agentes_procesar = {area_input: AGENTES[area_input]}
            else:
                print("El área ingresada no existe. Se procesarán todas las áreas.")
                agentes_procesar = AGENTES
        else:
            agentes_procesar = AGENTES

        respuestas_agentes = {}

        for nombre_agente, (carpeta, mensaje_sistema) in agentes_procesar.items():
            print(f"\n=== Procesando agente: {nombre_agente} ===")
            contenido_contexto, lista_archivos = leer_carpeta_recursiva(carpeta)
            prompt_agente = construir_prompt_agente(mensaje_sistema, contenido_contexto, lista_archivos, solicitud_usuario)
            respuesta = obtener_respuesta_api(prompt_agente)

            if respuesta:
                respuestas_agentes[nombre_agente] = respuesta
                print(f"\n--- Respuesta de {nombre_agente} ---\n")
                print(respuesta)
            else:
                respuestas_agentes[nombre_agente] = "Sin respuesta."
                print(f"No se obtuvo respuesta para {nombre_agente}.")

        # Si se procesaron más de 1 agente, se integra la respuesta del CEO
        if len(respuestas_agentes) > 1:
            print("\n=== Obteniendo la decisión final del 'CEO' ===")
            prompt_ceo = construir_prompt_ceo(respuestas_agentes)
            decision_ceo = obtener_respuesta_api(prompt_ceo)

            print("\n" + "="*60)
            if decision_ceo:
                print("DECISIÓN O CONCLUSIÓN FINAL DEL CEO:\n")
                print(decision_ceo)
            else:
                print("No se pudo obtener la respuesta final.")
            print("="*60)
        else:
            print("\nProcesamiento finalizado para el área seleccionada.")

# =====================================
# EJECUCIÓN DEL SCRIPT
# =====================================
if __name__ == "__main__":
    print("=== Sistema Multiagente con 10 agentes + CEO (integrador) ===")
    print("Soporta archivos PDF, DOC, DOCX, XLS, XLSX, CSV, TXT, PPTX.")
    print("Explora recursivamente las subcarpetas.")
    print("Instrucciones:")
    print("1. Crea 10 carpetas en el mismo directorio con estos nombres exactos:")
    print("   - contexto Sales")
    print("   - contexto Marketing")
    print("   - contexto Tools")
    print("   - contexto TasksPendientes")
    print("   - contexto Branding")
    print("   - contexto Finance")
    print("   - contexto LegalesDocumentos")
    print("   - contexto Clientes")
    print("   - contexto Web")
    print("   - contexto HR")
    print("2. Coloca en cada carpeta (o subcarpetas) los archivos relevantes.")
    print("3. Instala las dependencias: requests, PyPDF2, python-docx, pandas, openpyxl, python-pptx.")
    print("4. Ejecuta el script, ingresa tu solicitud y, opcionalmente, el área que deseas procesar.\n")
    main()



